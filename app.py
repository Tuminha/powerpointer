﻿from flask import Flask, render_template, request, send_from_directory, redirect, jsonify, send_file
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
import os
import openai
import collections.abc
from pptx import Presentation
from pptx.util import Inches
import random 
import re
import streamlit as st
import base64
import requests
from pptx.util import Pt
from dotenv import load_dotenv
from PIL import Image
from serp_images import main as serp_images
from io import BytesIO
from PIL import UnidentifiedImageError
from requests.exceptions import ConnectTimeout





from metapub import PubMedFetcher




load_dotenv()

# Load OpenAI API key from .env file
openai.api_key = os.getenv("OPENAI_API_KEY")

# Initialize the Flask app


app = Flask(__name__)

# Set a rate limit for the Flask app
limiter = Limiter(
    app,
    default_limits=["20 per day"], #This is the rate limit, you can remove it if you want
)

def get_pubmed_info(topic, start_year, end_year):
    fetch = PubMedFetcher()
    num_of_articles = 10

    # Add date range to the query
    query = f'{topic} AND {start_year}:{end_year}[Date - Publication]'

    # Get the PMID for the articles with the given topic
    pmids = fetch.pmids_for_query(query, retmax=num_of_articles)

    # Get articles
    articles_info = []
    for pmid in pmids:
        article = fetch.article_by_pmid(pmid)
        article_info = {
            "title": article.title,
            "abstract": article.abstract,
            "pub_date": article.year,
            "authors": ', '.join(article.authors)
        }
        articles_info.append(article_info)

    return articles_info




# Define the prompt used for the OpenAI API call

Prompt = """Write a detailed PowerPoint presentation about the user's topic, making use of the provided PubMed information. You only answer with the presentation. Follow the structure of the example.
Notice
- You do all the presentation text for the user.
- Each slide should contain at least 500 characters of content!
- Use the provided info from PubMed to add a short analysis in each slide when possible, highlighting potential biases, limitations, and other important information.
- You make the presentation easy to understand.
- The presentation starts with an introduction, covers specific subtopics, and ends with a conclusion.
- The presentation has a table of contents.
- The presentation has a summary.
- At least 10 slides.
- Each slide has a bibliographic reference.

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...
#Footer: Bibliographic reference

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Footer: Bibliographic reference

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Footer: Bibliographic reference

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE
#Footer: Bibliographic reference

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY
#Footer: Bibliographic reference

#Slide: END"""


# Function to generate text for the PPT using OpenAI API


def create_ppt_text(Input, start_year, end_year):
    # Get additional information from PubMed
    pubmed_info = get_pubmed_info(Input, start_year, end_year)

    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": (Prompt)},
            {"role": "user", "content": f"The user wants a presentation about {Input}. Here is some additional information I found on PubMed: {pubmed_info}"}
        ],
        temperature=0.5,
    )

    return response['choices'][0]['message']['content']


# Function to create PPT from the text generated by OpenAI API
def create_ppt(text_file, design_number, ppt_name, articles_info):
    # Create a presentation based on a design template
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    
    # Initialize variables
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True

    # Open the text file containing the presentation content
    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            # If the line starts with '#Title:', it's the title of the presentation
            if line.startswith('#Title:'):
                header = line.replace('#Title:', '').strip()
                # Add a new slide with the title layout (layout 0)
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                # Set the title of the slide
                title = slide.shapes.title
                title.text = header
                # Get the body shape of the slide
                body_shape = slide.shapes.placeholders[1]
                continue

            # If the line starts with '#Slide:', it's the start of a new slide
            elif line.startswith('#Slide:'):
                # If this is not the first slide, add the content to the previous slide
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content
                # Reset the content for the new slide
                content = "" 
                slide_count += 1

                # Select a random layout for the new slide, different from the last one
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]  # Possible layouts
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices) # Select random slide index
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue

            # If the line starts with '#Header:', it's the header for the slide
            elif line.startswith('#Header:'):
                header = line.replace('#Header:', '').strip()
                continue

            # If the line starts with '#Content:', it's the content for the slide
            elif line.startswith('#Content:'):
                content = line.replace('#Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

            # If the line starts with '#Footer:', it's the footer for the slide
            elif line.startswith('#Footer:'):
                footer = line.replace('#Footer:', '').strip()
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(6), Inches(0.5))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = footer
                p.font.size = Pt(12)
                continue

    # Add a references slide at the end
    slide_layout = prs.slide_layouts[5]  # Use the 'Title Slide' layout for the references slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "References"

    # Add a text box for the references
    left = Inches(0.5)
    top = Inches(1.0)
    width = height = Inches(6.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    for article_info in articles_info:
        p = tf.add_paragraph()
        p.text = f"{article_info['title']} by {article_info['authors']} ({article_info['pub_date']})"
        
    # Save the presentation
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"

    return file_path

# Function to generate download link for the created PPT

def get_ppt_download_link(ppt_file_path, ppt_name):
    with open(ppt_file_path, 'rb') as f:
        ppt_file = f.read()
    b64 = base64.b64encode(ppt_file).decode()
    href = f'<a href="data:file/ppt;base64,{b64}" download="{ppt_name}.pptx">Click here to download your presentation</a>'
    return href


# Main function to run the Streamlit app

def main():
    st.title('Periospot AI PowerPoint Generator')
    # Subheader
    st.subheader('Generate a PowerPoint presentation about any topic using the power of AI!')
    # Add a horizontal line
    st.markdown('---')
    # Paragraph
    st.write('This app uses the OpenAI API to generate a PowerPoint presentation about any topic. The presentation is generated based on the information available on PubMed about the topic.')    

    # Create a form block
    with st.form("my_form"):
        user_text = st.text_input('Enter your topic')
        design_number = st.selectbox('Choose a design:', options=range(1, 8), index=0)
        start_year = st.number_input('Start Year', min_value=1900, max_value=2100, value=2010, step=1)
        end_year = st.number_input('End Year', min_value=1900, max_value=2100, value=2023, step=1)
        
        # When the user presses the 'Submit' button, the form is submitted and all the inputs inside the form are used at once.
        submit_button = st.form_submit_button(label='Generate Presentation')

        # Provisional message
        link_placeholder = st.empty()

        # Spinners
        spinner_placeholder = st.empty()


        # Display provisional message
        link_placeholder.markdown("The link to your presentation will show up here.")

        if 'ppt_link' in st.session_state:
            st.markdown(st.session_state['ppt_link'], unsafe_allow_html=True)


        if submit_button and user_text:
            # Animate the text while the presentation is being generated
            spinner_placeholder.text('Generating presentation...')
            
            spinner_placeholder.empty() 

            # Clear previous system messages warnings and old links to the presentations
            st.session_state['ppt_link'] = None
            st.session_state['system_messages'] = None
            #Show a mesage to the user explaining that the presentation is being generated and all the logs will be displayed below. If some logs show error is not a problem, is just the way the API works.
            st.info('Generating presentation... Please wait. All the logs will be displayed below. If some logs show error is not a problem, is just the way the API works.')
            img_urls = serp_images(user_text)  
            for url in img_urls:
                try:
                    response = requests.get(url, timeout=5)
                    img = Image.open(BytesIO(response.content))
                    st.sidebar.image(img, caption=url, use_column_width=True)
                except UnidentifiedImageError:
                    st.error(f"Couldn't open image from URL: {url}")
                except ConnectTimeout:
                    st.warning(f"Timeout error when trying to connect to {url}. Skipping this image.")
                except requests.exceptions.RequestException as e:
                    st.error(f"An error occurred while trying to connect to {url}: {str(e)}")
            
            st.write("Design Number:", design_number, "selected.")
                    
            if design_number > 7:
                design_number = 1
                st.write("Unavailable design, using default design...")
            elif design_number == 0:
                design_number = 1
                st.write("Unavailable design, using default design...")
    
            with st.spinner('Fetching PubMed articles...'):
                articles_info = get_pubmed_info(user_text, start_year, end_year)
                for i, article_info in enumerate(articles_info, start=1):
                    st.write(f"Article {i}:")
                    st.write(f"Title: {article_info['title']}")
                    st.write(f"Publication date: {article_info['pub_date']}")
                    st.write(f"Authors: {article_info['authors']}")
                    st.write("---")
    
            st.info('PubMed articles fetched successfully. Now generating the PowerPoint presentation...')
    
            with open(f'Cache/{user_text}.txt', 'w', encoding='utf-8') as f:
                f.write(create_ppt_text(user_text, start_year, end_year))
    
            with st.spinner('Generating PowerPoint presentation...'):
                ppt_file_path = create_ppt(f'Cache/{user_text}.txt', design_number, user_text, articles_info)
    
            st.success('PowerPoint presentation generated successfully.')
            st.markdown(get_ppt_download_link(ppt_file_path, user_text), unsafe_allow_html=True)

            st.session_state['ppt_link'] = get_ppt_download_link(ppt_file_path, user_text)

            # Update the provisional message with the actual link
            if 'ppt_link' in st.session_state:
                link_placeholder.markdown(st.session_state['ppt_link'], unsafe_allow_html=True)

    


if __name__ == "__main__":
    main()


custom_css = """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Comfortaa:wght@300;700&display=swap');

    h1 {
        font-family: 'Comfortaa', sans-serif;
    }

    body {
        font-family: 'Comfortaa', sans-serif;
    }

    .stButton>button { 
        background-color: #1E90FF;
        color: white;
        font-family: 'Comfortaa', sans-serif;
    }

</style>
"""

st.markdown(custom_css, unsafe_allow_html=True)
